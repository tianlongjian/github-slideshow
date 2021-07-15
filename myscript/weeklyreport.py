#!/usr/bin/python
import csv
import xlrd
import xlwt
import re
#import test1
import pptx
import time
import datetime
from pptx import Presentation
import os
from xlutils.copy import copy as xl_copy


def getExcelData(file_name):
   
    list=[]
    a=[]
    work =xlrd.open_workbook(file_name)
    name=work.sheet_names()[0]
    print(name)
    sheet=work.sheet_by_name(name)
    print(sheet)
    nrows=sheet.nrows
    ncols=sheet.ncols
    print(ncols)
    print(nrows)
    #column=sheet.max_column
    #print(column)

    for i in range(nrows):
        print(i)
        list=[]
        if i==0:
            for j in range(ncols):
                if sheet.row_values(i)[j] == "项目":
                    project=j
                    print("~~~",project)
                    print(sheet.row_values(i)[project])   
                if sheet.row_values(i)[j] == "关键字":
                    keyword=j
                    print(keyword)
                    print(sheet.row_values(i)[keyword])  
                if sheet.row_values(i)[j] == "概要":
                    summ=j
                    print(summ)
                    print(sheet.row_values(i)[summ])
                if sheet.row_values(i)[j] == "状态":
                    status=j
                    print(status)
                    print(sheet.row_values(i)[status])
                if sheet.row_values(i)[j] == "已解决":
                    resovedtime=j
                    print(resovedtime)
                    print(sheet.row_values(i)[resovedtime])
        #print(sheet.row_values(i)[project])
        if len(sheet.row_values(i)[project])==0:
            print("no project")
            break
        if sheet.row_values(i)[project] is not None:
            list.append(sheet.row_values(i)[project])
        if sheet.row_values(i)[keyword] is not None:
            list.append(sheet.row_values(i)[keyword])
        if sheet.row_values(i)[summ] is not None:
            list.append(sheet.row_values(i)[summ])
        if sheet.row_values(i)[status] is not None:
            list.append(sheet.row_values(i)[status])
        if sheet.row_values(i)[resovedtime] is not None:
            list.append(sheet.row_values(i)[resovedtime])
        
        a.append(list)  
        
    #print(a)
    return a
def createEXCEL(data):
    if os.path.exists("bugreport.xlsx"):
        work =xlrd.open_workbook("bugreport.xlsx")
        book =xl_copy(work)
    else:
        book = xlwt.Workbook(encoding='utf-8',style_compression=0) 
    week=time.strftime("%W") 
    sheetname="W"+week
    #book.remove_sheet()
    sheet = book.add_sheet(sheetname,cell_overwrite_ok=True)
    print(len(data))
    print(len(data[0]))
    for i in range(len(data)):
        for j in range (len(data[i])):
            sheet.write(i,j,data[i][j] )
    #sheet.write(0,0,111)
    book.save('bugreport.xlsx' )
    #print(data)
    
    
def IsResovedBefore(excel_name,sheetname,inputkeyword):
    work =xlrd.open_workbook(excel_name)
    number= re.sub("W","",sheetname)
    name="W"+str(int(number)-1)
    print("IsResovedBefore",name)
    sheet=work.sheet_by_name(name)
    for i in range(sheet.nrows):
        if i == 0:
            for j in range(sheet.ncols):
                if sheet.row_values(i)[j] == "项目":
                    project=j
                    print("~~~",project)
                    print(sheet.row_values(i)[project])   
                if sheet.row_values(i)[j] == "关键字":
                    keyword=j
                    print(keyword)
                    print(sheet.row_values(i)[keyword])  
                if sheet.row_values(i)[j] == "概要":
                    summ=j
                    print(summ)
                    print(sheet.row_values(i)[summ])
                if sheet.row_values(i)[j] == "状态":
                    status=j
                    print(status)
                    print(sheet.row_values(i)[status])
                if sheet.row_values(i)[j] == "已解决":
                    resovedtime=j
                    print(resovedtime)
                    print(sheet.row_values(i)[resovedtime])
        else:
            print(sheet.row_values(i)[keyword],inputkeyword,sheet.row_values(i)[status])
            if sheet.row_values(i)[keyword]==inputkeyword:
                print("111111111111111111111111")
            if sheet.row_values(i)[status]=="已关闭":
                print("2222222222222222222222222")
            if sheet.row_values(i)[keyword]==inputkeyword and sheet.row_values(i)[status]=="已关闭":
                print("ture")
                return True
            
    return False            
    
    print("IsResovedBefore")
def getpptname():
    week2=time.strftime("%W")
    text=r"Engineering_Weekly_Report_[2021W11]_[WangShaowei]"
    text2=re.sub("W\d+\]","W"+week2+"]",text)
    name=text2+".pptx"
    return name

def WriteFirstPatePPT(pptname):
    prs= Presentation('template.pptx')
    slide = prs.slides.add_slide(prs.slide_layouts[0]) 
    #today1 = datetime.date.today()
    
    week2=time.strftime("%W")
   
    #week2=datetime.date.today().strptime('%W')
    print(week2)
    #print(today1)
    print(slide)
    print(slide.placeholders)
    for shape in slide.placeholders: 
        print(shape)
        phf = shape.placeholder_format
        text=r"Engineering_Weekly_Report_[2021W11]_[WangShaowei]"
        shape.text=re.sub("W\d+\]","W"+week2+"]",text)
        print(text)
        #pptname=shape.text+".pptx"
        print(f'{phf.idx}--{shape.name}--{phf.type}')
        #shape.text = f'{phf.idx}--{phf.type}'
    prs.save(pptname)
def WriteLastPatePPT(pptname):
    prs= Presentation(pptname)
    slide = prs.slides.add_slide(prs.slide_layouts[2]) 
    #today1 = datetime.date.today()
    week2=time.strftime("%W")
   
    #week2=datetime.date.today().strptime('%W')
    print(week2)
    #print(today1)
    print(slide)
    print(slide.placeholders)
    for shape in slide.placeholders: 
        print(shape)
        phf = shape.placeholder_format
        text=r"Engineering_Weekly_Report_[2021W11]_[WangShaowei]"
        shape.text=re.sub("W\d+\]","W"+week2+"]",text)
        print(text)
        print(f'{phf.idx}--{shape.name}--{phf.type}')
        #shape.text = f'{phf.idx}--{phf.type}'
    prs.save(pptname)
def writetopptfromExcel(excel_name,ppt_name): 
    list=[]
    a=[]
    work =xlrd.open_workbook(excel_name)
    week=time.strftime("%W") 
    name="W"+week
    #name=work.sheet_names()[0]
    print(name)
    sheet=work.sheet_by_name(name)
    print(sheet)
    nrows=sheet.nrows
    ncols=sheet.ncols
    print(ncols)
    print(nrows)
    WriteFirstPatePPT(ppt_name)
    print(pptname)
    for i in range(nrows):
        print(i)
        list=[]
        if i==0:
            for j in range(ncols):
                if sheet.row_values(i)[j] == "项目":
                    project=j
                    print("~~~",project)
                    print(sheet.row_values(i)[project])   
                if sheet.row_values(i)[j] == "关键字":
                    keyword=j
                    print(keyword)
                    print(sheet.row_values(i)[keyword])  
                if sheet.row_values(i)[j] == "概要":
                    summ=j
                    print(summ)
                    print(sheet.row_values(i)[summ])
                if sheet.row_values(i)[j] == "状态":
                    status=j
                    print(status)
                    print(sheet.row_values(i)[status])
                if sheet.row_values(i)[j] == "已解决":
                    resovedtime=j
                    print(resovedtime)
                    print(sheet.row_values(i)[resovedtime])
        else:
            if sheet.row_values(i)[status] == "已关闭" or sheet.row_values(i)[status] == "已解决" :
                title="DONE"
            else:
                title="ONGOING"
               
            ProjectName=sheet.row_values(i)[project]
            summary=sheet.row_values(i)[summ]
            bugnumber=sheet.row_values(i)[keyword]
            bugstatus=sheet.row_values(i)[status]    
            prs= Presentation(ppt_name)
            slide = prs.slides.add_slide(prs.slide_layouts[1]) 
            print("1",ProjectName,"2",summary,"3",bugnumber,"4",bugstatus)
            if bugstatus=="已关闭" and IsResovedBefore(excel_name,name,bugnumber):
                print("111111111111111111")
                continue
            else:
                
                for shape in slide.placeholders: 
                    print(shape)
                    phf = shape.placeholder_format
                    if phf.idx==0:
                        shape.text=title
                    if phf.idx==13:
                        shape.text=ProjectName+"\n"+bugnumber+"\n\n"+"问题描述：\n"+summary+"\n"+"问题分析："+"\n"+"问题解决："+"\n"+bugstatus
                    
                print(f'{phf.idx}--{shape.name}--{phf.type}')
                prs.save(ppt_name)
    WriteLastPatePPT(ppt_name)        
if __name__ == '__main__':
    file_name="mybug.xlsx"
    file_name2="bugreport.xlsx"
    pptname=getpptname()
    list =getExcelData(file_name)
    print(list)
    createEXCEL(list)
    #WriteFirstPatePPT('bugreport.pptx')
    writetopptfromExcel(file_name2,pptname)
   